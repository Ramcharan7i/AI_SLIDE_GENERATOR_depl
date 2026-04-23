import os
import random
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR


def create_ppt(slide_data, filename):
    try:
        # ✅ Ensure output folder exists
        os.makedirs("outputs/ppt", exist_ok=True)

        themes = [
            "templates/Berlin.pptx",
            "templates/Facet.pptx",
            "templates/Celestial.pptx",
            "templates/Depth.pptx",
            "templates/Mesh.pptx",
            "templates/Parallex.pptx",
            "templates/gallery.pptx",
            "templates/ion.pptx",
            "templates/BoardRoom.pptx",
            "templates/Circuit.pptx",
            "templates/Damask.pptx",
            "templates/Droplet.pptx",
            "templates/View.pptx",
            "templates/Slate.pptx",
            "templates/Wisp.pptx"
        ]

        template = random.choice(themes)

        prs = Presentation(template)

        # ✅ remove template slides safely
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        for i, slide in enumerate(slide_data):

            title = slide["title"].replace("*", "").strip()
            title = re.sub(r'^\d+[\.\)\-\s]*:', '', title)
            points = slide["points"][:4]
            image = slide.get("image")

            # -------- TITLE SLIDE --------
            if i == 0:
                layout = prs.slide_layouts[0]
                s = prs.slides.add_slide(layout)

                raw_title = slide["title"]

                match = re.search(r'(?i)about\s+(.*)', raw_title)
                clean_title = match.group(1) if match else raw_title

                clean_title = re.sub(r'(?i)\b(explain|explaining|presentation|slides)\b', '', clean_title)
                clean_title = re.sub(r'^\d+[\.\)\-\s]*', '', clean_title).strip()

                if not clean_title:
                    clean_title = raw_title

                s.shapes.title.text = clean_title.title()
                s.placeholders[1].text = "LearnLift - AI"

            # -------- IMAGE + TEXT --------
            elif image and os.path.exists(image):
                layout = prs.slide_layouts[5]
                s = prs.slides.add_slide(layout)

                s.shapes.title.text = title

                top_position = Inches(2.3)

                textbox = s.shapes.add_textbox(
                    Inches(0.7),
                    top_position,
                    Inches(5),
                    Inches(3.8)
                )

                tf = textbox.text_frame
                tf.word_wrap = True
                tf.clear()

                for j, p in enumerate(points[:5]):
                    para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    para.text = p.strip()
                    para.font.size = Pt(24)

                # ✅ safe image add
                try:
                    s.shapes.add_picture(
                        image,
                        Inches(6.3),
                        top_position,
                        height=Inches(3.8)
                    )
                except Exception as e:
                    print("IMAGE ADD ERROR:", e)

            # -------- NORMAL SLIDE --------
            else:
                layout = prs.slide_layouts[1]
                s = prs.slides.add_slide(layout)

                s.shapes.title.text = title
                tf = s.shapes.placeholders[1].text_frame

                for j, p in enumerate(points):
                    para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    para.text = p
                    para.font.size = Pt(24)

            # ✅ title font
            s.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)

        # ✅ FINAL SAVE PATH (FIXED)
        output_path = f"outputs/ppt/{filename}.pptx"

        prs.save(output_path)

        return output_path

    except Exception as e:
        print("PPT GENERATION ERROR:", e)
        return None